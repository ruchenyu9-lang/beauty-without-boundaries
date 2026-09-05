"""Generate a simple project presentation PPT for 美妆无界."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──────────────────────────────────────────
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
BG_GRAD    = RGBColor(0x66, 0x7E, 0xEA)   # purple-blue
ACCENT     = RGBColor(0xC4, 0x45, 0x69)   # rose
ACCENT2    = RGBColor(0xD4, 0x73, 0x5E)   # coral
ACCENT3    = RGBColor(0x4A, 0x2C, 0x2A)   # dark brown
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    """Add solid background rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_accent_bar(slide, left, top, width=Inches(1.2), height=Pt(4), color=ACCENT):
    """Add a thin colored accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_feature_card(slide, left, top, icon, title, subtitle, desc, color):
    """Add a feature card with icon, title, and description."""
    card_w = Inches(3.6)
    card_h = Inches(4.2)
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Icon
    add_text(slide, left + Inches(0.2), top + Inches(0.2), Inches(1), Inches(0.8), icon, size=36, color=color, bold=True)
    # Title
    add_text(slide, left + Inches(0.2), top + Inches(0.9), card_w - Inches(0.4), Inches(0.5), title, size=20, color=WHITE, bold=True)
    # Subtitle
    add_text(slide, left + Inches(0.2), top + Inches(1.35), card_w - Inches(0.4), Inches(0.4), subtitle, size=13, color=color)
    # Accent bar
    add_accent_bar(slide, left + Inches(0.2), top + Inches(1.75), Inches(0.8), Pt(3), color)
    # Description
    add_text(slide, left + Inches(0.2), top + Inches(1.9), card_w - Inches(0.4), Inches(2.0), desc, size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════
# Slide 1: Title
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# Top accent strip
strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
         "美妆无界", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.6),
         "Beauty Without Boundaries", size=22, color=ACCENT, bold=False, align=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(5.5), Inches(3.7), Inches(2.3), Pt(4), ACCENT)

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
         "欧莱雅第二届美妆黑客松 · 无界体验家赛道", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.6),
         "多模态美妆交互平台 — 让美妆科技真正包容每个人", size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# Slide 2: Problem & Vision
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
         "痛点与愿景", size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.5), Pt(4), ACCENT)

# Pain points
pains = [
    "🔹  色弱/色盲用户无法准确辨别妆容色号，试妆困难",
    "🔹  上传妆容图后缺少专业分析与步骤指导",
    "🔹  线上试妆缺乏个性化色号推荐，色号匹配全凭感觉",
]
for i, p in enumerate(pains):
    add_text(slide, Inches(1.0), Inches(1.5 + i * 0.65), Inches(5.5), Inches(0.6),
             p, size=15, color=LIGHT_GRAY)

# Vision
add_text(slide, Inches(7.2), Inches(0.5), Inches(5.5), Inches(0.6),
         "我们的解法", size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(7.2), Inches(1.15), Inches(1.5), Pt(4), ACCENT)

visions = [
    "💡  AI 驱动的智能色号匹配与妆容方案生成",
    "💡  上传即分析：AI 步骤拆解 + 平替产品推荐",
    "💡  Daltonize 色彩增强 + 色盲检测 + 定制可感知调色板",
]
for i, v in enumerate(visions):
    add_text(slide, Inches(7.4), Inches(1.5 + i * 0.65), Inches(5.5), Inches(0.6),
             v, size=15, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# Slide 3: 3 Core Features
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
         "3 大核心功能", size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.05), Inches(1.5), Pt(4), ACCENT)

features = [
    ("💄", "智能妆容生成", "3D扫描 + AI搭配",
     "录入已有产品色号，3D扫描人脸，AI智能生成多种妆容方案并标记使用色号。实时人脸追踪叠加妆容效果。", ACCENT),
    ("📸", "妆容图片分析", "AI分析 + 平替推荐",
     "上传妆容图片，AI自动分析化妆步骤和注意事项，智能推荐平替产品，降低决策成本。", ACCENT2),
    ("🎨", "色觉无障碍", "色彩增强 + 色彩翻译",
     "Daltonize色彩增强算法、Ishihara色盲检测、色彩翻译描述、为色觉异常用户定制可感知调色板。", ACCENT3),
]

gap = Inches(0.35)
card_w = Inches(3.6)
total_w = card_w * 3 + gap * 2
start_left = (W - total_w) // 2

for i, (icon, title, subtitle, desc, color) in enumerate(features):
    left = start_left + i * (card_w + gap)
    add_feature_card(slide, left, Inches(1.5), icon, title, subtitle, desc, color)

# ══════════════════════════════════════════════════════════
# Slide 4: Tech Stack
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
         "技术栈", size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.5), Pt(4), ACCENT)

tech_groups = [
    ("前端", ["React + TypeScript", "Vite 构建工具", "MediaPipe Face Mesh 人脸追踪", "Web Speech API 语音合成"], BG_GRAD),
    ("后端", ["Python + FastAPI", "SQLite + SQLAlchemy ORM", "Uvicorn ASGI 服务器", "Mock 数据 (无需外部 API)"], ACCENT2),
    ("色彩科学", ["CIEDE2000 ΔE 色差算法", "CIE Lab 色彩空间转换", "Daltonize 色彩增强", "Brettel/Viénot 模拟矩阵"], ACCENT3),
]

for gi, (group_name, items, color) in enumerate(tech_groups):
    left = Inches(0.8 + gi * 4.0)
    top = Inches(1.6)

    # Group card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.6), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_text(slide, left + Inches(0.3), top + Inches(0.2), Inches(3), Inches(0.5),
             group_name, size=22, color=color, bold=True)
    add_accent_bar(slide, left + Inches(0.3), top + Inches(0.75), Inches(0.8), Pt(3), color)

    for j, item in enumerate(items):
        add_text(slide, left + Inches(0.3), top + Inches(1.0 + j * 0.7), Inches(3), Inches(0.6),
                 f"▸  {item}", size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════
# Slide 5: Architecture
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
         "系统架构", size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.5), Pt(4), ACCENT)

# Architecture boxes
boxes = [
    (Inches(0.8), Inches(1.8), Inches(3.5), Inches(2.5), "前端 React + Vite",
     ["• SmartMakeupPage — 人脸追踪+妆容叠加", "• ImageAnalysisPage — 图片上传+分析", "• ColorVisionPage — 色觉检测+调色板",
      "• 全局 ColorVisionContext 状态管理"], BG_GRAD),
    (Inches(5.0), Inches(1.8), Inches(3.5), Inches(2.5), "后端 FastAPI",
     ["• /api/makeup — 妆容生成+色号匹配", "• /api/analysis — 图片分析+平替", "• /api/color-vision — 色觉增强+检测",
      "• /api/products — 产品数据服务"], ACCENT2),
    (Inches(9.2), Inches(1.8), Inches(3.5), Inches(2.5), "数据层",
     ["• SQLite + SQLAlchemy ORM", "• products.json / shades.json 种子数据",
      "• MockClaudeService 预生成结果", "• 无需外部 AI API 密钥"], ACCENT3),
]

for (left, top, w, h, title, items, color) in boxes:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_text(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.4),
             title, size=16, color=color, bold=True)

    for j, item in enumerate(items):
        add_text(slide, left + Inches(0.2), top + Inches(0.6 + j * 0.42), w - Inches(0.4), Inches(0.4),
                 item, size=11, color=LIGHT_GRAY)

# Arrows (simple text arrows between boxes)
add_text(slide, Inches(4.35), Inches(2.8), Inches(0.6), Inches(0.4), "⟶", size=24, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.55), Inches(2.8), Inches(0.6), Inches(0.4), "⟶", size=24, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

# Flow summary
add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(1.2),
         "用户浏览器 ⟶  React 前端 (Vite Dev / Nginx Prod)  ⟶  FastAPI 后端  ⟶  SQLite + 种子数据\n"
         "色觉增强算法 (Daltonize / CIEDE2000 / Lab) 完全在前端 Client-Side 运行，无需后端计算",
         size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# Slide 6: API Endpoints
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
         "API 端点一览", size=32, color=WHITE, bold=True)
add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.5), Pt(4), ACCENT)

endpoints = [
    ("妆容生成", BG_GRAD, [
        ("POST", "/api/makeup/scan-face", "3D人脸扫描"),
        ("POST", "/api/makeup/shade-match", "色号匹配"),
        ("POST", "/api/makeup/generate-looks", "妆容方案生成"),
        ("POST", "/api/makeup/render-overlay-data", "妆容叠加数据"),
    ]),
    ("图片分析", ACCENT2, [
        ("POST", "/api/analysis/analyze-image", "AI图片分析"),
        ("POST", "/api/analysis/find-substitutes", "平替产品搜索"),
    ]),
    ("色觉无障碍", ACCENT3, [
        ("POST", "/api/color-vision/daltonize-single", "Daltonize色彩增强"),
        ("POST", "/api/color-vision/detect-type", "色盲类型检测"),
        ("GET",  "/api/color-vision/ishihara-plates", "Ishihara测试图"),
        ("POST", "/api/color-vision/generate-palette", "定制调色板生成"),
    ]),
]

col = 0
for group_name, color, eps in endpoints:
    left = Inches(0.8 + col * 4.0)
    top = Inches(1.6)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.6), Inches(0.5 + len(eps) * 0.65))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x45)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_text(slide, left + Inches(0.2), top + Inches(0.1), Inches(3), Inches(0.4),
             group_name, size=16, color=color, bold=True)

    for j, (method, path, desc) in enumerate(eps):
        y = top + Inches(0.55 + j * 0.65)
        method_color = RGBColor(0x6B, 0xD4, 0x7A) if method == "GET" else RGBColor(0x6B, 0x9E, 0xF0)
        add_text(slide, left + Inches(0.2), y, Inches(0.6), Inches(0.3),
                 method, size=10, color=method_color, bold=True)
        add_text(slide, left + Inches(0.8), y, Inches(2.7), Inches(0.3),
                 path, size=10, color=LIGHT_GRAY)
        add_text(slide, left + Inches(0.8), y + Inches(0.22), Inches(2.7), Inches(0.3),
                 desc, size=9, color=RGBColor(0x88, 0x88, 0x88))

    col += 1

# ══════════════════════════════════════════════════════════
# Slide 7: Demo / Thank You
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT
strip.line.fill.background()

add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
         "谢谢观看", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(5.5), Inches(3.3), Inches(2.3), Pt(4), ACCENT)

add_text(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
         "Beauty Without Boundaries — 让美妆科技真正包容每个人", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         "🚀  在线体验:  http://localhost:5173", size=14, color=GOLD, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
         "欧莱雅第二届美妆黑客松 · 无界体验家赛道", size=12, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────
output_path = r"D:\claudeworkspace\beauty-without-boundaries\美妆无界_项目展示.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
